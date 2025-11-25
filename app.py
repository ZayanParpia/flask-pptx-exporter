"""
Flask PPTX exporter (final)
- Click template thumbnails (static/images) to choose template
- Preview lists all slide pairs (Top / Bottom) and is fully editable
- On export we create slides and remove the original first (template) slide
- Top text centered, Bottom text left-aligned
- 'Quran' option removed from available templates/images
- Important fix: bottom lines are written as separate paragraphs,
  leading TABs/control chars and leading normal spaces removed,
  paragraph indents/levels explicitly reset to avoid accidental extra indentation.
  Zeroed text-frame margins to remove template padding.
  Marsiya-specific footer and watermark added.
- NEW: bottom phrase is reflowed so every 5 words starts a new line (keeps formatting).
"""

import os
import re
import uuid
import tempfile
from flask import Flask, render_template, request, send_file, abort
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# new import for safe filenames
from werkzeug.utils import secure_filename

# configure paths - adjust to your environment
TEMPLATES_FOLDER = r"C:\Users\ikepa\OneDrive\Pictures\JCC\templates"
STATIC_FOLDER = r"C:\Users\ikepa\OneDrive\Pictures\JCC\static"

app = Flask(__name__, template_folder=TEMPLATES_FOLDER, static_folder=STATIC_FOLDER)
@app.route("/tutorial", methods=["GET"])
def tutorial_route():
    return render_template("tutorial.html")


# formatting hints (optional)
FORMATS = {
    "marsiya": {
        "Top":    {"color": "#ffc000", "font": "Open Sans", "size": 54, "bold": True},
        "Bottom": {"color": "#ffffff", "font": "Open Sans", "size": 40, "bold": True},
    },
    "qasida": {
        "Top":    {"color": "#e4c0a8", "font": "Arial", "size": 54, "bold": True},
        "Bottom": {"color": "#ffffff", "font": "Calibri", "size": 44, "bold": True},
    }
}

# -----------------------
# Helpers
# -----------------------
def hex_to_rgb_tuple(hex_color):
    if not hex_color:
        return (0, 0, 0)
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        return (0, 0, 0)
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def ensure_templates():
    files = []
    if not os.path.isdir(TEMPLATES_FOLDER):
        return files
    for f in os.listdir(TEMPLATES_FOLDER):
        if f.lower().endswith(".pptx"):
            files.append(f)
    return sorted(files)

def find_matching_template_for_image(image_basename, pptx_list):
    base = image_basename.lower()
    for p in pptx_list:
        if base in os.path.splitext(p.lower())[0]:
            return p
    candidate = base.capitalize() + ".pptx"
    if candidate in pptx_list:
        return candidate
    return None

def copy_font_attrs(src_run, dst_run):
    """Copy basic font attributes from src_run to dst_run if present."""
    try:
        font = src_run.font
        if getattr(font, "name", None):
            dst_run.font.name = font.name
        if getattr(font, "size", None):
            dst_run.font.size = font.size
        if getattr(font, "bold", None) is not None:
            dst_run.font.bold = font.bold
        if getattr(font, "italic", None) is not None:
            dst_run.font.italic = font.italic
        try:
            if font.color and getattr(font.color, "rgb", None):
                dst_run.font.color.rgb = font.color.rgb
        except Exception:
            pass
    except Exception:
        pass

def apply_fmt_to_run(run, fmt):
    """Legacy: unconditional apply of fmt. Kept for compatibility but not used below."""
    if not fmt:
        return
    if fmt.get("font"):
        run.font.name = fmt["font"]
    if fmt.get("size"):
        try:
            run.font.size = Pt(int(fmt["size"]))
        except Exception:
            pass
    if fmt.get("bold") is not None:
        run.font.bold = bool(fmt["bold"])
    if fmt.get("italic") is not None:
        run.font.italic = bool(fmt["italic"])
    if fmt.get("color"):
        r, g, b = hex_to_rgb_tuple(fmt["color"])
        run.font.color.rgb = RGBColor(r, g, b)

def apply_fmt_respecting_template(run, fmt):
    """
    Apply fmt to run, but do NOT override attributes already present on the run (template wins).
    Only sets properties that are missing/None on the run's font.
    """
    if not fmt:
        return
    font = run.font
    # font name: apply only if template doesn't provide one
    try:
        if fmt.get("font") and not getattr(font, "name", None):
            font.name = fmt["font"]
    except Exception:
        pass

    # size: apply only if not present
    try:
        if fmt.get("size") and getattr(font, "size", None) is None:
            font.size = Pt(int(fmt["size"]))
    except Exception:
        pass

    # bold: apply only if template left it unset (None)
    try:
        if fmt.get("bold") is not None and getattr(font, "bold", None) is None:
            font.bold = bool(fmt["bold"])
    except Exception:
        pass

    # italic
    try:
        if fmt.get("italic") is not None and getattr(font, "italic", None) is None:
            font.italic = bool(fmt["italic"])
    except Exception:
        pass

    # color: only set if there is no existing rgb color
    try:
        if fmt.get("color"):
            existing_rgb = None
            if getattr(font, "color", None) and getattr(font.color, "rgb", None):
                existing_rgb = font.color.rgb
            if existing_rgb is None:
                r, g, b = hex_to_rgb_tuple(fmt["color"])
                font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass

def set_paragraph_alignment(paragraph, align):
    """Set paragraph alignment; accepts PP_ALIGN or strings."""
    try:
        if isinstance(align, str):
            a = align.lower()
            if a == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif a == "left":
                paragraph.alignment = PP_ALIGN.LEFT
            elif a == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = align
    except Exception:
        pass

def reset_paragraph_format(paragraph):
    """
    Explicitly clear paragraph-level indentation and spacing that often
    remain from templates and cause extra indentation on subsequent lines.
    """
    try:
        pf = paragraph.paragraph_format
        # zero out indents and spacing
        pf.left_indent = Pt(0)
        pf.first_line_indent = Pt(0)
        pf.space_before = Pt(0)
        pf.space_after = Pt(0)
        # ensure paragraph is top-level (no bullet/level-based indent)
        try:
            paragraph.level = 0
        except Exception:
            pass
        # attempt to remove bullet/level remnants (best-effort)
        try:
            # no-op tolerant call to private API region; safe to ignore failures
            paragraph._p.get_or_add_pPr()
        except Exception:
            pass
    except Exception:
        pass

def sanitize_line(line, strip_leading_spaces=True):
    """
    Remove problematic characters that can produce tabs/indentation in the PPTX.
    Removes control chars, explicit tabs, NBSPs and zero-width spaces.
    Optionally strips normal leading spaces to prevent Google Slides collapsing lines.
    """
    if line is None:
        return ""
    s = line
    # remove explicit tabs
    s = s.replace("\t", "")
    # remove carriage returns (we split on \n elsewhere)
    s = s.replace("\r", "")
    # normalize NBSP to normal space
    s = s.replace("\u00A0", " ")
    # remove zero-width and BOM characters
    s = s.replace("\u200b", "").replace("\ufeff", "")
    # remove remaining ASCII control characters (except newline if present)
    s = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", "", s)
    # remove leading tab-like / NBSP / ZWSP leftover at start
    s = re.sub(r"^[\t\u00A0\u200b\ufeff]+", "", s)
    # optionally remove normal leading spaces too (Google Slides collapses them)
    if strip_leading_spaces:
        s = s.lstrip()
    # trim trailing newline artifacts but keep normal spaces
    s = s.rstrip("\n")
    return s

# -----------------------
# NEW helper: break every N words
# -----------------------
def break_every_n_words(text, n=5):
    """
    Break a text into lines where each line contains up to n words.
    Returns a single string with '\n' line separators.
    - Keeps punctuation attached to words.
    - Collapses existing whitespace and joins across original newlines.
    """
    if not text:
        return ""
    # take tokens separated by whitespace
    words = re.findall(r"\S+", text)
    if not words:
        return ""
    lines = []
    for i in range(0, len(words), n):
        lines.append(" ".join(words[i:i + n]))
    return "\n".join(lines)

# -----------------------
# Normalization helper (programmatic Slide Master edit)
# -----------------------
def normalize_template_placeholders(prs):
    """
    Walk slide masters, layouts and the first slide to normalize any placeholder
    text frames so template-level indents/margins/bullet levels do not carry over.
    This is a best-effort programmatic substitute for 'Edit Slide Master'.
    """
    def _normalize_shape(sh):
        if not hasattr(sh, "text_frame") or sh.text_frame is None:
            return
        try:
            tf = sh.text_frame
            # Zero margins on placeholder text frames
            try:
                tf.margin_left = Inches(0)
                tf.margin_right = Inches(0)
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
            except Exception:
                pass

            # Ensure paragraphs are reset and runs sanitized
            try:
                for p in tf.paragraphs:
                    reset_paragraph_format(p)
                    for r in p.runs:
                        try:
                            r.text = sanitize_line(r.text)
                        except Exception:
                            pass
            except Exception:
                pass
        except Exception:
            pass

    # Normalize slide master placeholders
    try:
        sm = prs.slide_master
        for sh in sm.shapes:
            _normalize_shape(sh)
    except Exception:
        pass

    # Normalize layouts
    try:
        for layout in prs.slide_layouts:
            for sh in layout.shapes:
                _normalize_shape(sh)
    except Exception:
        pass

    # Also normalize the actual first slide if present (common templates)
    try:
        if len(prs.slides) > 0:
            for sh in prs.slides[0].shapes:
                _normalize_shape(sh)
    except Exception:
        pass

# -----------------------
# Helpers to write each line as paragraphs or explicit breaks
# -----------------------
def write_lines_to_textframe_preserve_style(text_frame, text, src_run_sample=None, override_fmt=None, align=None, use_paragraphs=False):
    """
    Write text into text_frame.
    - If use_paragraphs is False: uses single paragraph + explicit run breaks (original behavior).
    - If use_paragraphs is True: writes all lines into the text_frame.text using '\n'.join,
      letting python-pptx turn each '\n' into a new paragraph.
    Leading spaces are stripped (by default) to avoid Google Slides collapsing lines.
    """
    # Clear existing content first
    text_frame.clear()

    # Zero text-frame margins to avoid template-inherited padding/indentation.
    try:
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)
    except Exception:
        pass

    # Build lines list - preserve empty lines as empty visual lines
    lines = (text or "").splitlines()
    if len(lines) == 0:
        lines = [""]

    # sanitize lines
    lines = [sanitize_line(ln, strip_leading_spaces=True) for ln in lines]

    if use_paragraphs:
        # Simple approach: set the whole text to newline-joined sanitized lines.
        # python-pptx will convert each '\n' to a separate paragraph.
        try:
            text_frame.text = "\n".join(lines)
        except Exception:
            # fallback - manually add paragraphs
            text_frame.clear()
            for i, ln in enumerate(lines):
                if i == 0:
                    p = text_frame.paragraphs[0]
                    p.text = ln or ""
                else:
                    p = text_frame.add_paragraph()
                    p.text = ln or ""
    else:
        # Single paragraph with explicit breaks between runs (original behavior)
        p = text_frame.paragraphs[0]
        reset_paragraph_format(p)
        if align is not None:
            set_paragraph_alignment(p, align)

        prev_run = None
        for raw_line in lines:
            line = sanitize_line(raw_line, strip_leading_spaces=True)
            run = p.add_run()
            if prev_run is not None:
                try:
                    prev_run.add_break()
                except Exception:
                    try:
                        prev_run.text = (prev_run.text or "") + "\n"
                    except Exception:
                        pass
            run.text = line or ""
            # preserve template attributes if present
            if src_run_sample:
                try:
                    copy_font_attrs(src_run_sample, run)
                except Exception:
                    pass
            # apply override_fmt only where template didn't provide the attribute
            try:
                apply_fmt_respecting_template(run, override_fmt)
            except Exception:
                pass
            prev_run = run

    # After writing, attempt to normalize paragraph formatting and apply formatting
    try:
        for p in text_frame.paragraphs:
            reset_paragraph_format(p)
            if align is not None:
                set_paragraph_alignment(p, align)
            for r in p.runs:
                # first, copy template sample run attributes if we had a sample but run lacks attrs
                try:
                    if src_run_sample:
                        # If run has no name/size/bold/etc, copy from sample (copy_font_attrs copies only present attrs on sample)
                        copy_font_attrs(src_run_sample, r)
                except Exception:
                    pass
                # Then apply override_fmt only where template didn't provide them
                try:
                    apply_fmt_respecting_template(r, override_fmt)
                except Exception:
                    pass
    except Exception:
        pass

def set_text_preserve_shape(shape, text, override_fmt=None, align=None, use_paragraphs=False):
    """
    For an existing shape: sample the first run (if present) to preserve its style,
    ensure zeroed margins, then write lines using write_lines_to_textframe_preserve_style.
    """
    tf = shape.text_frame
    if tf is None:
        return
    sample_run = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        sample_run = tf.paragraphs[0].runs[0]

    # Zero text-frame margins to remove accidental indentation/padding from template placeholders
    try:
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
    except Exception:
        pass

    write_lines_to_textframe_preserve_style(tf, text, src_run_sample=sample_run, override_fmt=override_fmt, align=align, use_paragraphs=use_paragraphs)

def add_marsiya_footer_and_watermark(slide, prs):
    """Add centered footer and a small red '.' watermark at bottom-right for Marsiya template."""
    try:
        # Footer across full width -> center aligned
        footer_box = slide.shapes.add_textbox(Inches(0), prs.slide_height - Inches(0.6), prs.slide_width, Inches(0.5))
        tf = footer_box.text_frame
        tf.clear()
        # zero margins
        try:
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)
        except Exception:
            pass

        p = tf.paragraphs[0]
        reset_paragraph_format(p)
        set_paragraph_alignment(p, PP_ALIGN.CENTER)
        r = p.add_run()
        r.text = "Marshia Translations | ISIJ of Toronto"
        # apply Calibri 14 bold white
        try:
            r.font.name = "Calibri"
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255, 255, 255)
        except Exception:
            pass

        # small red dot on bottom-right as watermark
        dot_size = Inches(0.25)
        dot_left = prs.slide_width - dot_size - Inches(0.15)
        dot_top = prs.slide_height - dot_size - Inches(0.15)
        dot_box = slide.shapes.add_textbox(dot_left, dot_top, dot_size, dot_size)
        tf2 = dot_box.text_frame
        tf2.clear()
        try:
            tf2.margin_left = Inches(0)
            tf2.margin_right = Inches(0)
            tf2.margin_top = Inches(0)
            tf2.margin_bottom = Inches(0)
        except Exception:
            pass
        p2 = tf2.paragraphs[0]
        reset_paragraph_format(p2)
        set_paragraph_alignment(p2, PP_ALIGN.RIGHT)
        r2 = p2.add_run()
        r2.text = "."
        try:
            r2.font.size = Pt(28)
            r2.font.bold = True
            r2.font.color.rgb = RGBColor(255, 0, 0)
        except Exception:
            pass
    except Exception:
        # non-fatal; don't break generation if watermark/footer can't be added
        pass

# -----------------------
# Core generator
# -----------------------
def generate_pptx_from_texts(template_filename, top_block, bottom_block):
    template_path = os.path.join(TEMPLATES_FOLDER, template_filename)
    if not os.path.exists(template_path):
        raise FileNotFoundError("Template not found: " + template_path)

    # Keep empty lines (preview maps 1:1)
    top_lines = [ln.rstrip("\r") for ln in (top_block or "").splitlines()]
    bottom_lines = [ln.rstrip("\r") for ln in (bottom_block or "").splitlines()]

    max_len = max(len(top_lines), len(bottom_lines))
    if max_len == 0:
        top_lines = [""]
        bottom_lines = [""]
        max_len = 1

    # normalize lengths
    while len(top_lines) < max_len:
        top_lines.append("")
    while len(bottom_lines) < max_len:
        bottom_lines.append("")

    # Load template presentation
    prs = Presentation(template_path)

    # NEW: normalize placeholders/layouts/slide-master so indent/bullets don't leak in
    try:
        normalize_template_placeholders(prs)
    except Exception:
        # non-fatal: continue even if normalization fails
        pass

    # Choose a prototype layout (first slide's layout if available)
    proto_slide = prs.slides[0] if len(prs.slides) > 0 else None
    proto_layout = proto_slide.slide_layout if proto_slide is not None else prs.slide_layouts[6]

    def get_text_shapes_for_proto(slide):
        shapes = []
        if not slide:
            return shapes
        for sh in slide.shapes:
            if hasattr(sh, "text_frame") and sh.text_frame is not None:
                shapes.append(sh)
        try:
            shapes.sort(key=lambda s: s.top)
        except Exception:
            pass
        return shapes

    proto_text_shapes = get_text_shapes_for_proto(proto_slide)
    use_placeholders = len(proto_text_shapes) >= 2

    # determine format hints
    base_name = os.path.splitext(template_filename)[0].lower()
    fmt_key = None
    for key in FORMATS:
        if key.lower() == base_name:
            fmt_key = key
            break
    template_formats = FORMATS.get(fmt_key, {})

    # Create slides for each pair
    for idx in range(max_len):
        new_slide = prs.slides.add_slide(proto_layout)
        top_text = top_lines[idx]
        bottom_text = bottom_lines[idx]

        # NEW: reflow bottom_text so every 5 words are on a new line
        # Combine if bottom_text had embedded newlines (we treat entire phrase)
        try:
            joined_bottom = " ".join([ln for ln in bottom_text.splitlines() if ln.strip() != ""])
            bottom_text_reflowed = break_every_n_words(joined_bottom, n=5)
        except Exception:
            bottom_text_reflowed = bottom_text or ""

        if use_placeholders:
            new_text_shapes = get_text_shapes_for_proto(new_slide)
            if len(new_text_shapes) >= 2:
                top_shape = new_text_shapes[0]
                bottom_shape = new_text_shapes[-1]
                # ensure frames have zero margins to avoid leftover indent/padding
                try:
                    if top_shape.text_frame:
                        top_shape.text_frame.margin_left = Inches(0)
                        top_shape.text_frame.margin_right = Inches(0)
                        top_shape.text_frame.margin_top = Inches(0)
                        top_shape.text_frame.margin_bottom = Inches(0)
                except Exception:
                    pass
                try:
                    if bottom_shape.text_frame:
                        bottom_shape.text_frame.margin_left = Inches(0)
                        bottom_shape.text_frame.margin_right = Inches(0)
                        bottom_shape.text_frame.margin_top = Inches(0)
                        bottom_shape.text_frame.margin_bottom = Inches(0)
                except Exception:
                    pass

                # Top: preserve prior single-paragraph behavior (centered)
                # Note: we prefer template attributes (font/size/color/position). override_fmt will only fill missing attrs.
                set_text_preserve_shape(top_shape, top_text, override_fmt=template_formats.get("Top"), align=PP_ALIGN.CENTER, use_paragraphs=False)
                # Bottom: write the reflowed bottom text (every 5 words -> new line) as paragraphs
                set_text_preserve_shape(bottom_shape, bottom_text_reflowed, override_fmt=template_formats.get("Bottom"), align=PP_ALIGN.LEFT, use_paragraphs=True)

                # If this is marsiya template, add footer + watermark
                try:
                    if base_name == "marsiya":
                        add_marsiya_footer_and_watermark(new_slide, prs)
                except Exception:
                    pass

                continue

        # fallback textboxes (no placeholders found) - these will use code positioning
        left = Inches(0.7)
        top_pos = Inches(0.6)
        width = prs.slide_width - Inches(1.4)

        top_box = new_slide.shapes.add_textbox(left, top_pos, width, Inches(2.5))
        # zero margins on textframe
        try:
            top_box.text_frame.margin_left = Inches(0)
            top_box.text_frame.margin_right = Inches(0)
            top_box.text_frame.margin_top = Inches(0)
            top_box.text_frame.margin_bottom = Inches(0)
        except Exception:
            pass
        # write lines (top - single paragraph behavior)
        write_lines_to_textframe_preserve_style(top_box.text_frame, top_text, src_run_sample=None, override_fmt=template_formats.get("Top"), align=PP_ALIGN.CENTER, use_paragraphs=False)

        bottom_top = prs.slide_height - Inches(2.8)
        bottom_box = new_slide.shapes.add_textbox(left, bottom_top, width, Inches(2.5))
        try:
            bottom_box.text_frame.margin_left = Inches(0)
            bottom_box.text_frame.margin_right = Inches(0)
            bottom_box.text_frame.margin_top = Inches(0)
            bottom_box.text_frame.margin_bottom = Inches(0)
        except Exception:
            pass
        

        # write bottom - paragraph per line (works with Google Slides) using reflowed bottom text
        write_lines_to_textframe_preserve_style(bottom_box.text_frame, bottom_text_reflowed, src_run_sample=None, override_fmt=template_formats.get("Bottom"), align=PP_ALIGN.LEFT, use_paragraphs=True)

        # If this is marsiya template, add footer + watermark
        try:
            if base_name == "marsiya":
                add_marsiya_footer_and_watermark(new_slide, prs)
        except Exception:
            pass

    # Remove original template first slide so exported file doesn't include the template slide
    try:
        sldIdLst = prs.slides._sldIdLst  # private API but commonly used
        if len(sldIdLst) > 0:
            sldIdLst.remove(sldIdLst[0])
    except Exception:
        # non-fatal - continue to save
        pass

    out_path = os.path.join(tempfile.gettempdir(), f"export_{uuid.uuid4().hex}.pptx")
    prs.save(out_path)
    return out_path

# -----------------------
# Filename sanitization helper
# -----------------------
def make_safe_pptx_filename(raw_name: str, fallback_base: str = "export", max_len: int = 120) -> str:
    """
    Turn a user-provided name into a safe filename. Uses werkzeug.secure_filename
    to strip dangerous characters, ensures .pptx extension, and enforces max length.
    If the result is empty, returns fallback_base + '_export.pptx'.
    """
    if not raw_name:
        base = fallback_base
    else:
        base = raw_name.strip()

    # remove any path separators just in case
    base = base.replace("/", " ").replace("\\", " ")

    # use secure_filename to remove unsafe characters
    safe = secure_filename(base)

    # fallback if secure_filename produced empty string
    if not safe:
        safe = fallback_base

    # ensure extension
    name, ext = os.path.splitext(safe)
    if ext.lower() != ".pptx":
        safe = f"{name}.pptx"

    # enforce max length (leave room for .pptx)
    if len(safe) > max_len:
        name_only = os.path.splitext(safe)[0][: max_len - 5]
        safe = f"{name_only}.pptx"

    return safe

# -----------------------
# Flask routes
# -----------------------
@app.route("/", methods=["GET"])
def index_route():
    # available pptx templates (exclude quran)
    all_pptx = ensure_templates()
    available_templates = [p for p in all_pptx if os.path.splitext(p)[0].lower() != "quran"]

    # gather images from static/images but exclude quran image
    images = []
    images_dir = os.path.join(STATIC_FOLDER, "images")
    if os.path.isdir(images_dir):
        for f in sorted(os.listdir(images_dir)):
            if f.lower().endswith((".png", ".jpg", ".jpeg", ".webp")) and os.path.splitext(f)[0].lower() != "quran":
                images.append(f)

    # build a mapping image -> matched template filename (if found)
    image_to_template = {}
    for img in images:
        base = os.path.splitext(img)[0]
        matched = find_matching_template_for_image(base, available_templates)
        image_to_template[img] = matched

    return render_template("index.html", templates=available_templates, images=images, image_to_template=image_to_template)

@app.route("/generate", methods=["POST"])
def generate_route():
    template_file = request.form.get("template")
    top_text = request.form.get("top_text", "")
    bottom_text = request.form.get("bottom_text", "")
    # read custom filename from the form (HTML input name="pptx_name")
    raw_name = request.form.get("pptx_name", "").strip()

    if not template_file:
        abort(400, "No template selected")

    # security: only allow filenames that exist in templates folder (and exclude quran)
    available = ensure_templates()
    available = [p for p in available if os.path.splitext(p)[0].lower() != "quran"]

    if template_file not in available:
        abort(400, "Template not available")

    try:
        out_path = generate_pptx_from_texts(template_file, top_text, bottom_text)
    except Exception as e:
        abort(500, f"Generation error: {e}")

    # make a safe download filename from user input; default to template-based name if blank/invalid
    fallback_base = os.path.splitext(template_file)[0] + "_export"
    safe_download_name = make_safe_pptx_filename(raw_name or fallback_base, fallback_base=fallback_base)

    # send file with user-chosen safe name
    try:
        return send_file(out_path, as_attachment=True, download_name=safe_download_name)
    except TypeError:
        # older Flask versions use 'attachment_filename'
        return send_file(out_path, as_attachment=True, attachment_filename=safe_download_name)

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)

